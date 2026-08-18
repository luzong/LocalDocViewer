# -*- coding: utf-8 -*-
"""
本地文档书籍资料浏览器 - 数据扫描与缩略图提取器 (v3)

新特性:
 1. GUI 界面, 可勾选扫描的文件类型v3
 2. 支持自定义提取图片数量 (1-10 张), 按文件实际页数裁剪
 3. 扩展支持: PPT/PPTX, EPUB, HTML/CSS/JS/Python/CPP/Java/Go 等编程语言
 4. 基于持久化哈希缓存, 实现增量扫描:
    - 文件未变动 (mtime + size 一致) → 跳过
    - 文件新增/变动 → 生成缩略图
    - 磁盘已删除文件 → 从缓存清除, 对应缩略图删除
 5. 命令行与 GUI 两种使用模式

输出:
  - site/assets/data/data.json
  - site/data.js
  - site/assets/thumbnails/*.png
  - site/assets/data/scan_cache.json (增量缓存)
"""

import os
import sys
import json
import time
import hashlib
import traceback
import threading
import shutil
import ctypes
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Set, Dict, List, Optional, Tuple


# ================================================================
# 应用图标 (多尺寸 icon.ico) 资源路径
#   - 桌面/资源管理器图标: 由 PyInstaller --icon 嵌入 exe
#   - 窗口左上角图标:      self.root.iconbitmap(icon_path)
#   - 任务栏/Alt-Tab 图标: SetCurrentProcessExplicitAppUserModelID (需在 tk.Tk() 前调用)
# ================================================================
APP_USER_MODEL_ID = "DocumentScanner.LocalBrowser.v1"  # Win7+ 任务栏分组 ID
ICON_PATH = Path(__file__).resolve().parent / "icon.ico"


def set_app_user_model_id(app_id: str = APP_USER_MODEL_ID) -> None:
    """在 Windows 上设置进程的 AppUserModelID, 让任务栏/Alt-Tab 显示本应用图标而非默认 Python 图标.

    必须在创建主窗口 (Tk) 之前调用, 否则不生效.
    """
    if sys.platform != "win32":
        return
    try:
        ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID(app_id)
    except Exception:
        # 非 Windows 或 API 不可用, 忽略
        pass


def get_icon_path() -> Optional[str]:
    """返回 icon.ico 的绝对路径字符串, 不存在则返回 None.

    优先使用与 scanner.py 同目录的 icon.ico;
    若是 PyInstaller 打包后运行, 则解压到 sys._MEIPASS 临时目录查找.
    """
    candidates = []
    if hasattr(sys, "_MEIPASS"):
        candidates.append(Path(sys._MEIPASS) / "icon.ico")
    candidates.append(Path(__file__).resolve().parent / "icon.ico")
    for p in candidates:
        if p.is_file():
            return str(p)
    return None

# ---------------- 核心依赖检查 ----------------
MISSING = []
try:
    import fitz  # PyMuPDF
except ImportError:
    MISSING.append("PyMuPDF")

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    MISSING.append("Pillow")

if MISSING:
    print("缺少核心依赖: " + ", ".join(MISSING))
    print("请运行: pip install -r requirements.txt")
    sys.exit(1)

# 可选依赖
try:
    import docx  # python-docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import win32com.client  # pywin32
    HAS_WIN32 = True
except ImportError:
    HAS_WIN32 = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from ebooklib import epub
    HAS_EPUB = True
except ImportError:
    HAS_EPUB = False


# ================================================================
# 支持的扩展名与类型分组
# ================================================================
# 类型分组 (供 GUI 显示)
TYPE_GROUPS = {
    "Documents": {
        "label": "文档 (Documents)",
        "exts": ["pdf", "doc", "docx", "rtf"],
    },
    "Slides": {
        "label": "幻灯片 (Slides)",
        "exts": ["ppt", "pptx"],
    },
    "Ebooks": {
        "label": "电子书 (Ebooks)",
        "exts": ["epub", "mobi", "azw3"],
    },
    "Text": {
        "label": "文本 (Text)",
        "exts": ["txt", "md", "rtf", "log", "csv", "json", "xml", "yaml", "yml", "ini", "cfg", "conf", "bat", "sh", "ps1"],
    },
    "Code_Web": {
        "label": "Web 前端",
        "exts": ["html", "htm", "css", "js", "jsx", "ts", "tsx", "vue", "svelte"],
    },
    "Code_Script": {
        "label": "脚本语言",
        "exts": ["py", "rb", "pl", "php", "bat", "sh", "ps1", "r", "jl", "lua"],
    },
    "Code_Compile": {
        "label": "编译/系统语言",
        "exts": ["c", "h", "cpp", "cxx", "cc", "hpp", "hh", "cs", "java", "go", "rs", "kt", "swift", "d", "f", "sql"],
    },
    "Code_DotNet": {
        "label": ".NET / 微软技术栈",
        "exts": ["cs", "vb", "aspx", "cshtml", "razor"],
    },
}

# 扩展名 -> 分类 (与前端一致)
SUPPORTED_EXTS: Dict[str, str] = {}
# 文本类 (显示代码预览)
TEXT_CATEGORIES = {"txt", "md", "rtf", "log", "csv", "json", "xml", "yaml", "yml",
                   "ini", "cfg", "conf", "bat", "sh", "ps1", "r", "jl", "lua",
                   "html", "htm", "css", "js", "jsx", "ts", "tsx", "vue", "svelte",
                   "py", "rb", "pl", "php", "c", "h", "cpp", "cxx", "cc", "hpp", "hh",
                   "cs", "java", "go", "rs", "kt", "swift", "d", "f", "sql",
                   "vb", "aspx", "cshtml", "razor"}

def _build_ext_map():
    """根据 TYPE_GROUPS 构建扩展名 -> 分类 映射"""
    _EBOOK_EXTS = {"epub", "mobi", "azw3"}
    for group_key, group in TYPE_GROUPS.items():
        for ext in group["exts"]:
            cat = group_key  # 用分组名作为 category
            if ext in TEXT_CATEGORIES:
                cat = "text"
            elif ext in ("pdf",):
                cat = "pdf"
            elif ext in ("doc", "docx"):
                cat = "word"
            elif ext in ("ppt", "pptx"):
                cat = "ppt"
            elif ext in _EBOOK_EXTS:
                cat = "ebook"
            else:
                cat = "text"
            SUPPORTED_EXTS[ext] = cat

_build_ext_map()

# 去重后的扩展名列表 (rtf 同时在 Documents 和 Text, 优先 Documents 的分类)
# 重新调整 rtF
SUPPORTED_EXTS["rtf"] = "text"  # RTF 按文本处理, 便于解析

# 缩略图尺寸
THUMB_W = 400
THUMB_H = 560   # 纵向, 像书的封面比例

# 文本缩略图最多行数与字符
TXT_MAX_LINES = 22
TXT_MAX_CHARS = 60

# 跳过目录
SKIP_DIRS = {".git", ".svn", "__pycache__", "node_modules", "$RECYCLE.BIN", "System Volume Information"}


# ================================================================
# 工具
# ================================================================
def safe_id(path_str: str) -> str:
    h = hashlib.md5(path_str.encode("utf-8")).hexdigest()
    return h[:12]


def human_size(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if n < 1024:
            return f"{n:.1f} {unit}"
        n /= 1024
    return f"{n:.1f} PB"


def ensure_dir(p: Path):
    p.mkdir(parents=True, exist_ok=True)


def file_signature(path: str) -> Tuple[int, float]:
    """返回 (size, mtime) 作为文件变动签名"""
    try:
        st = os.stat(path)
        return st.st_size, st.st_mtime
    except Exception:
        return -1, -1.0


def _default_font(size: int = 22):
    candidates = [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for f in candidates:
        if os.path.exists(f):
            try:
                return ImageFont.truetype(f, size)
            except Exception:
                pass
    return ImageFont.load_default()


def _truncate(s: str, n: int) -> str:
    s = (s or "").replace("\t", " ").replace("\r", " ")
    if len(s) <= n:
        return s
    return s[:n - 1] + "…"


def _draw_text_thumbnail(title: str, subtitle: str, lines: list, out_path: Path,
                         bg_color=(245, 247, 250), border_color=(70, 110, 180),
                         title_color=(30, 40, 60), subtitle_color=(110, 120, 140),
                         text_color=(50, 60, 80)):
    img = Image.new("RGB", (THUMB_W, THUMB_H), bg_color)
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, THUMB_W, 70], fill=border_color)
    title_font = _default_font(26)
    d.text((20, 18), _truncate(title, 22), font=title_font, fill=(255, 255, 255))
    sub_font = _default_font(16)
    d.text((20, 80), _truncate(subtitle, 40), font=sub_font, fill=subtitle_color)
    text_font = _default_font(18)
    y = 120
    for ln in lines[:TXT_MAX_LINES]:
        d.text((20, y), _truncate(ln, TXT_MAX_CHARS), font=text_font, fill=text_color)
        y += 24
    foot_font = _default_font(14)
    d.text((20, THUMB_H - 30), "Local Document Browser", font=foot_font, fill=(170, 180, 200))
    img.save(out_path, "PNG")


# ================================================================
# 持久化缓存
# ================================================================
@dataclass
class CacheEntry:
    path: str              # 绝对路径
    size: int              # 文件大小
    mtime: float           # 修改时间
    thumb_id: str          # 缩略图文件名(不含路径)
    ext: str = ""
    category: str = ""

    def signature(self) -> Tuple[int, float]:
        return self.size, self.mtime

    def matches(self, new_size: int, new_mtime: float) -> bool:
        return self.size == new_size and abs(self.mtime - new_mtime) < 0.1


class FileCache:
    """基于 scan_cache.json 的增量扫描缓存"""

    def __init__(self, cache_path: Path):
        self.path = cache_path
        self.entries: Dict[str, CacheEntry] = {}
        self.load()

    def load(self):
        if self.path.exists():
            try:
                with open(self.path, "r", encoding="utf-8") as f:
                    raw = json.load(f)
                for p, v in raw.items():
                    self.entries[p] = CacheEntry(
                        path=p, size=v["size"], mtime=v["mtime"],
                        thumb_id=v["thumb_id"], ext=v.get("ext", ""),
                        category=v.get("category", "")
                    )
            except Exception as e:
                print(f"  [warn] 缓存文件损坏, 将重建: {e}")
                self.entries = {}

    def save(self):
        self.path.parent.mkdir(parents=True, exist_ok=True)
        raw = {}
        for p, e in self.entries.items():
            raw[p] = {
                "size": e.size, "mtime": e.mtime,
                "thumb_id": e.thumb_id, "ext": e.ext, "category": e.category,
            }
        tmp = self.path.with_suffix(".tmp")
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(raw, f, ensure_ascii=False, indent=2)
        tmp.replace(self.path)

    def is_changed(self, path: str) -> bool:
        """检查文件是否变动或新增"""
        if path not in self.entries:
            return True  # 新文件
        ent = self.entries[path]
        size, mtime = file_signature(path)
        return not ent.matches(size, mtime)

    def is_processed(self, path: str) -> bool:
        return path in self.entries and not self.is_changed(path)

    def get(self, path: str) -> Optional[CacheEntry]:
        return self.entries.get(path)

    def put(self, path: str, thumb_id: str, ext: str, category: str):
        size, mtime = file_signature(path)
        self.entries[path] = CacheEntry(
            path=path, size=size, mtime=mtime,
            thumb_id=thumb_id, ext=ext, category=category
        )

    def remove(self, path: str):
        self.entries.pop(path, None)

    def all_paths(self) -> Set[str]:
        return set(self.entries.keys())

    def sync_deleted(self, current_paths: Set[str], thumb_dir: Path, log=print):
        """删除已不存在的文件的缓存项与缩略图"""
        deleted = set()
        for p in self.all_paths():
            if p not in current_paths:
                deleted.add(p)

        if not deleted:
            return 0

        for p in deleted:
            ent = self.entries.get(p)
            if ent:
                thumb_path = thumb_dir / ent.thumb_id
                if thumb_path.exists():
                    try:
                        thumb_path.unlink()
                        log(f"  [rm] 已删除: {p}")
                    except Exception:
                        pass
                # 清理多页缩略图 (_2.png, _3.png ...)
                stem = thumb_path.stem
                for i in range(2, 11):
                    multi_path = thumb_dir / f"{stem}_{i}.png"
                    if multi_path.exists():
                        try:
                            multi_path.unlink()
                        except Exception:
                            pass
                    else:
                        break
            self.remove(p)
        return len(deleted)


# ================================================================
# 缩略图生成 (支持多页)
# ================================================================
def gen_pdf_thumbnail(pdf_path: str, out_path: Path, page_count: int = 1) -> bool:
    """生成 PDF 缩略图: 第一页作为封面 out_path, 多页时额外保存 _2, _3..."""
    try:
        doc = fitz.open(pdf_path)
        if doc.page_count == 0:
            doc.close()
            return False
        pages_to_render = min(page_count, doc.page_count)
        # 每页都用完整尺寸渲染, 不再拼接压缩
        zoom = THUMB_W / 72.0 * 1.3
        mat = fitz.Matrix(zoom, zoom)
        for i in range(pages_to_render):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            # 缩放到统一尺寸
            img = img.resize((THUMB_W, THUMB_H), Image.LANCZOS)
            if i == 0:
                img.save(out_path, "PNG", optimize=True)
            else:
                # 多页: 保存为 out_path_2.png, out_path_3.png ...
                multi_path = out_path.with_name(
                    out_path.stem + f"_{i+1}" + out_path.suffix
                )
                img.save(multi_path, "PNG", optimize=True)
        doc.close()
        return True
    except Exception as e:
        print(f"  [warn] PDF 缩略图失败 {pdf_path}: {e}")
        return False


def gen_word_thumbnail(word_path: str, ext: str, out_path: Path, page_count: int = 1) -> bool:
    if HAS_WIN32:
        tmp_pdf = out_path.with_suffix(".tmp.pdf")
        if _word_to_pdf_via_win32(word_path, str(tmp_pdf)):
            ok = gen_pdf_thumbnail(str(tmp_pdf), out_path, page_count=page_count)
            try:
                tmp_pdf.unlink(missing_ok=True)
            except Exception:
                pass
            if ok:
                return True
    # 降级: 读取文本
    text_lines = []
    title = os.path.basename(word_path)
    if ext == "docx" and HAS_DOCX:
        try:
            d = docx.Document(word_path)
            for p in d.paragraphs:
                t = p.text.strip()
                if t:
                    text_lines.append(t)
                if len(text_lines) >= TXT_MAX_LINES:
                    break
        except Exception as e:
            print(f"  [warn] docx 读取失败 {word_path}: {e}")
    else:
        text_lines = ["(无法解析此 Word 文档)", "请使用 Office 打开查看内容"]
    _draw_text_thumbnail(title, f"Word Document  .{ext}", text_lines, out_path,
                         border_color=(43, 87, 154))
    return True


def _word_to_pdf_via_win32(word_path: str, pdf_path: str) -> bool:
    if not HAS_WIN32:
        return False
    word = None
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        doc = word.Documents.Open(os.path.abspath(word_path), ReadOnly=True)
        doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17)
        doc.Close(False)
        return os.path.exists(pdf_path)
    except Exception as e:
        print(f"  [warn] Word COM 转换失败 {word_path}: {e}")
        return False
    finally:
        if word:
            try:
                word.Quit()
            except Exception:
                pass


def gen_ppt_thumbnail(ppt_path: str, ext: str, out_path: Path, page_count: int = 1) -> bool:
    """生成 PPT 缩略图: 优先转 PDF 渲染, 否则提取幻灯片文字"""
    # 1) 尝试用 COM 转 PDF
    if HAS_WIN32:
        tmp_pdf = out_path.with_suffix(".tmp.pdf")
        if _ppt_to_pdf_via_win32(ppt_path, str(tmp_pdf)):
            ok = gen_pdf_thumbnail(str(tmp_pdf), out_path, page_count=page_count)
            try:
                tmp_pdf.unlink(missing_ok=True)
            except Exception:
                pass
            if ok:
                return True

    # 2) python-pptx 提取文本
    text_lines = []
    title = os.path.basename(ppt_path)
    if HAS_PPTX and ext == "pptx":
        try:
            prs = Presentation(ppt_path)
            max_slides = min(page_count, len(prs.slides))
            for i in range(max_slides):
                slide = prs.slides[i]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in para.runs).strip()
                            if t:
                                text_lines.append(t)
                            if len(text_lines) >= TXT_MAX_LINES:
                                break
                        if len(text_lines) >= TXT_MAX_LINES:
                            break
                if len(text_lines) >= TXT_MAX_LINES:
                    break
        except Exception as e:
            print(f"  [warn] pptx 读取失败 {ppt_path}: {e}")
    else:
        text_lines = [f"({ext.upper()} 文件, 建议安装 Office 生成真实预览)"]

    if not text_lines:
        text_lines = ["(空文件或读取失败)"]
    _draw_text_thumbnail(title, f"PowerPoint  .{ext}", text_lines, out_path,
                         border_color=(207, 76, 50))
    return True


def _ppt_to_pdf_via_win32(ppt_path: str, pdf_path: str) -> bool:
    if not HAS_WIN32:
        return False
    app = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(os.path.abspath(ppt_path), ReadOnly=True, WithWindow=False)
        pres.SaveAs(os.path.abspath(pdf_path), FileFormat=32)  # 32 = wdFormatPDF
        pres.Close()
        return os.path.exists(pdf_path)
    except Exception as e:
        print(f"  [warn] PPT COM 转换失败 {ppt_path}: {e}")
        return False
    finally:
        if app:
            try:
                app.Quit()
            except Exception:
                pass


def gen_epub_thumbnail(epub_path: str, out_path: Path) -> bool:
    """电子书 (EPUB/MOBI/AZW3): 提取封面或首页文本, 失败则显示通用缩略图"""
    title = os.path.basename(epub_path)
    ext = os.path.splitext(epub_path)[1].lower().lstrip(".")
    text_lines = []

    if ext == "epub" and HAS_EPUB:
        try:
            book = epub.read_epub(epub_path)
            try:
                if book.get_metadata("DC", "title"):
                    title = book.get_metadata("DC", "title")[0][0]
            except Exception:
                pass
            items = list(book.get_items_of_type(9))  # 9 = ITEM_DOCUMENT
            for item in items[:3]:
                try:
                    content = item.get_content().decode("utf-8", errors="ignore")
                    import re
                    clean = re.sub(r"<[^>]+>", " ", content)
                    for line in clean.split("\n"):
                        line = line.strip()
                        if line and len(line) > 3:
                            text_lines.append(line)
                        if len(text_lines) >= TXT_MAX_LINES:
                            break
                    if len(text_lines) >= TXT_MAX_LINES:
                        break
                except Exception:
                    continue
        except Exception as e:
            print(f"  [warn] epub 解析失败 {epub_path}: {e}")

    if not text_lines:
        if ext == "epub":
            text_lines = ["(EPUB 文件预览)", "请使用 EPUB 阅读器打开"]
        elif ext == "mobi":
            text_lines = ["(MOBI 电子书)", "Kindle 格式, 请使用 Kindle 阅读器打开"]
        elif ext == "azw3":
            text_lines = ["(AZW3 电子书)", "Kindle 格式, 请使用 Kindle 阅读器打开"]
        else:
            text_lines = [f"({ext.upper()} 电子书)", "请使用对应阅读器打开"]

    _draw_text_thumbnail(title, f"Ebook  .{ext}", text_lines, out_path,
                         border_color=(120, 80, 150))
    return True


def gen_text_thumbnail(file_path: str, out_path: Path, ext: str = "") -> bool:
    """生成文本/代码文件缩略图"""
    title = os.path.basename(file_path)
    lines = []
    try:
        for enc in ("utf-8", "gbk", "gb18030", "utf-16"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    for line in f:
                        line = line.rstrip("\n").rstrip("\r")
                        if line.strip():
                            lines.append(line)
                        if len(lines) >= TXT_MAX_LINES:
                            break
                break
            except UnicodeDecodeError:
                continue
            except Exception:
                break
    except Exception as e:
        print(f"  [warn] 文件读取失败 {file_path}: {e}")
    if not lines:
        lines = ["(空文件或读取失败)"]
    border_map = {
        "py": (53, 114, 176), "js": (237, 189, 37), "ts": (49, 120, 198),
        "html": (227, 76, 38), "css": (38, 77, 228), "java": (231, 81, 39),
        "cpp": (69, 105, 180), "c": (0, 102, 140), "go": (0, 173, 216),
        "rs": (255, 90, 60), "cs": (139, 69, 19), "vb": (56, 129, 203),
        "php": (119, 63, 143), "rb": (204, 52, 52), "lua": (0, 0, 128),
        "bat": (90, 90, 90), "sh": (60, 130, 50),
    }
    color = border_map.get(ext, (60, 140, 90))
    _draw_text_thumbnail(title, f"{'Code' if ext in TEXT_CATEGORIES else 'Text'}  .{ext or 'txt'}", lines, out_path,
                         border_color=color)
    return True


def gen_generic_thumbnail(file_path: str, ext: str, out_path: Path) -> bool:
    title = os.path.basename(file_path)
    info_lines = [
        f"文件类型: .{ext}",
        f"大小: {human_size(os.path.getsize(file_path) if os.path.exists(file_path) else 0)}",
        "",
        "暂不支持生成内容预览",
        "请使用对应软件打开",
    ]
    _draw_text_thumbnail(title, f"Document  .{ext}", info_lines, out_path,
                         border_color=(120, 100, 130))
    return True


# ================================================================
# Web 服务器 (静态文件 + 本地操作 API)
# ================================================================
def _guess_mime(path: str) -> str:
    """根据扩展名推断 MIME 类型, 供浏览器内联预览使用"""
    import mimetypes
    mtype, _ = mimetypes.guess_type(path)
    if mtype:
        return mtype
    ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""
    # mimetypes 库覆盖不全的常见类型
    map_extra = {
        "pdf": "application/pdf",
        "md": "text/markdown; charset=utf-8",
        "py": "text/plain; charset=utf-8",
        "js": "application/javascript; charset=utf-8",
        "css": "text/css; charset=utf-8",
        "json": "application/json; charset=utf-8",
        "yaml": "text/yaml; charset=utf-8",
        "yml": "text/yaml; charset=utf-8",
        "csv": "text/csv; charset=utf-8",
        "log": "text/plain; charset=utf-8",
    }
    if ext in map_extra:
        return map_extra[ext]
    return "application/octet-stream"


class WebServer:
    """基于 http.server 的本地 Web 服务器, 支持静态文件和 API 路由"""

    def __init__(self, site_dir: str, port: int = 8080, log_fn=None):
        import http.server
        import socketserver
        self.http_server = http.server
        self.socketserver = socketserver
        self.site_dir = str(site_dir)
        self.port = port
        self.log_fn = log_fn or print
        self._server = None
        self._thread = None

        # 动态创建 Handler 类 (闭包捕获 self)
        srv = self

        class Handler(http.server.SimpleHTTPRequestHandler):
            def __init__(self, *args, **kwargs):
                super().__init__(*args, directory=srv.site_dir, **kwargs)

            def _send_json(self, data, code=200):
                body = json.dumps(data, ensure_ascii=False).encode("utf-8")
                self.send_response(code)
                self.send_header("Content-Type", "application/json; charset=utf-8")
                self.send_header("Content-Length", len(body))
                self.end_headers()
                self.wfile.write(body)

            def do_GET(self):
                """静态文件 + API GET"""
                if self.path.startswith("/api/"):
                    return self._handle_api_get()
                return super().do_GET()

            def do_POST(self):
                """API POST 路由"""
                if self.path.startswith("/api/"):
                    return self._handle_api_post()
                self.send_error(404, "Not Found")

            def _handle_api_get(self):
                if self.path == "/api/status":
                    return self._send_json({"ok": True, "service": "doc-browser"})
                if self.path.startswith("/api/serve-file"):
                    return self._serve_file()
                self.send_error(404, "Unknown API")

            def _serve_file(self):
                """通过服务器代理文件让浏览器直接预览 (避免 file:// 被拦截)"""
                from urllib.parse import urlparse, unquote, quote
                parsed = urlparse(self.path)
                params = {}
                if parsed.query:
                    for kv in parsed.query.split("&"):
                        if "=" in kv:
                            k, v = kv.split("=", 1)
                            params[k] = unquote(v)
                path = params.get("path", "").strip()
                if not path or not os.path.exists(path):
                    self.send_error(404, "File not found")
                    return
                if not os.path.isfile(path):
                    self.send_error(400, "Not a file")
                    return
                try:
                    size = os.path.getsize(path)
                    ctype = _guess_mime(path)
                    self.send_response(200)
                    self.send_header("Content-Type", ctype)
                    self.send_header("Content-Length", size)
                    fname = os.path.basename(path)
                    self.send_header("Content-Disposition",
                                     f"inline; filename*=UTF-8''{quote(fname)}")
                    self.send_header("Cache-Control", "public, max-age=3600")
                    self.end_headers()
                    with open(path, "rb") as f:
                        while True:
                            chunk = f.read(1024 * 256)
                            if not chunk:
                                break
                            self.wfile.write(chunk)
                except BrokenPipeError:
                    pass
                except Exception as e:
                    try:
                        self.send_error(500, f"Internal error: {e}")
                    except Exception:
                        pass

            def _handle_api_post(self):
                content_len = int(self.headers.get("Content-Length", 0))
                raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                try:
                    body = json.loads(raw.decode("utf-8"))
                except Exception:
                    return self._send_json({"ok": False, "error": "请求格式错误"}, 400)

                path = body.get("path", "").strip()
                if not path:
                    return self._send_json({"ok": False, "error": "缺少 path 参数"}, 400)

                if self.path == "/api/open-file":
                    return self._open_file(path)
                elif self.path == "/api/open-folder":
                    return self._open_folder(path)
                elif self.path == "/api/copy-path":
                    return self._copy_path(path)
                else:
                    self._send_json({"ok": False, "error": "未知 API"}, 404)

            def _open_file(self, path):
                try:
                    if not os.path.exists(path):
                        return self._send_json({"ok": False, "error": "文件不存在"}, 404)
                    if os.name == "nt":
                        os.startfile(path)
                    elif sys.platform == "darwin":
                        import subprocess
                        subprocess.Popen(["open", path])
                    else:
                        import subprocess
                        subprocess.Popen(["xdg-open", path])
                    return self._send_json({"ok": True})
                except Exception as e:
                    return self._send_json({"ok": False, "error": str(e)}, 500)

            def _open_folder(self, path):
                try:
                    if not os.path.exists(path):
                        return self._send_json({"ok": False, "error": "文件不存在"}, 404)
                    if os.name == "nt":
                        # explorer /select,"文件路径" 选中该文件
                        import subprocess
                        subprocess.Popen(f'explorer /select,"{path}"', shell=True)
                    elif sys.platform == "darwin":
                        import subprocess
                        subprocess.Popen(["open", "-R", path])
                    else:
                        import subprocess
                        folder = os.path.dirname(path) or "."
                        subprocess.Popen(["xdg-open", folder])
                    return self._send_json({"ok": True})
                except Exception as e:
                    return self._send_json({"ok": False, "error": str(e)}, 500)

            def _copy_path(self, path):
                try:
                    import subprocess
                    if os.name == "nt":
                        subprocess.Popen(f'echo|set /p="{path}"|clip', shell=True)
                    else:
                        import pyperclip
                        pyperclip.copy(path)
                    return self._send_json({"ok": True})
                except Exception as e:
                    return self._send_json({"ok": False, "error": str(e)}, 500)

            def log_message(self, fmt, *args):
                # 静默或通过 log_fn 输出
                pass

        self.Handler = Handler

    def start(self):
        """启动服务器 (非阻塞, 在后台线程运行)"""
        if self._server:
            return False, "服务器已在运行"
        try:
            self._server = self.socketserver.ThreadingTCPServer(
                ("127.0.0.1", self.port), self.Handler
            )
            self._server.daemon_threads = True
            self._thread = threading.Thread(target=self._server.serve_forever, daemon=True)
            self._thread.start()
            self.log_fn(f"✓ Web 服务已启动: http://127.0.0.1:{self.port}")
            return True, f"http://127.0.0.1:{self.port}"
        except OSError as e:
            if e.errno == 98 or "Address already in use" in str(e):
                # 端口被占, 尝试下一个
                self.log_fn(f"端口 {self.port} 被占用, 尝试 {self.port+1}")
                self.port += 1
                return self.start()
            self.log_fn(f"[error] 启动失败: {e}")
            return False, str(e)

    def stop(self):
        """停止服务器"""
        if not self._server:
            return
        try:
            self._server.shutdown()
            self._server.server_close()
            self._server = None
            self._thread = None
            self.log_fn("Web 服务已停止")
        except Exception as e:
            self.log_fn(f"[error] 停止失败: {e}")

    @property
    def running(self):
        return self._server is not None

    @property
    def url(self):
        return f"http://127.0.0.1:{self.port}"


# ================================================================
# 扫描主流程
# ================================================================
class Scanner:
    def __init__(self, root: str, output_dir: str,
                 ext_filter: Set[str] = None,
                 image_count: int = 1,
                 force_rescan: bool = False,
                 log_fn=None):
        self.root = os.path.abspath(root)
        self.output_dir = Path(output_dir).absolute()
        self.thumb_dir = self.output_dir / "assets" / "thumbnails"
        self.data_dir = self.output_dir / "assets" / "data"
        ensure_dir(self.thumb_dir)
        ensure_dir(self.data_dir)
        self.root_data_json = self.output_dir / "data.json"

        # 过滤的扩展名集合, None 表示全部
        self.ext_filter = ext_filter if ext_filter else set(SUPPORTED_EXTS.keys())
        # 去除重复
        self.ext_filter = set(e.lower().lstrip(".") for e in self.ext_filter)

        self.image_count = max(1, min(image_count, 10))
        self.force_rescan = force_rescan

        self.log_fn = log_fn or print
        self.files_flat = []
        self.current_paths: Set[str] = set()
        self.thumb_count = 0
        self.fail_count = 0
        self.skip_count = 0
        self.skipped_ext = 0
        self.new_count = 0
        self.cached_count = 0

        # 进度回调
        self.progress_cb = None

        # 停止标志 (由外部线程设置, 扫描循环中检查)
        self._cancelled = False

        # 进度统计 (扫描中动态更新)
        self._total_files = 0
        self._processed_files = 0

        # 持久化缓存
        self.cache = FileCache(self.data_dir / "scan_cache.json")

    def _ensure_site_structure(self):
        """检测输出目录是否有网站结构, 没有则从模板复制生成"""
        index_html = self.output_dir / "index.html"
        if index_html.exists():
            return  # 已有网站结构

        # 查找网站模板目录 (优先级: scanner/data > site > site_template)
        scanner_dir = Path(__file__).parent
        candidates = [
            scanner_dir / "data",
            scanner_dir.parent / "site",
            scanner_dir / "site_template",
        ]
        template_dir = None
        for c in candidates:
            if (c / "index.html").exists():
                template_dir = c
                break

        if not template_dir:
            self.log("⚠ 未找到网站模板, 无法自动生成网站结构")
            return

        self.log(f"输出目录缺少网站结构, 从模板复制: {template_dir}")
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # 复制模板文件 (index.html, css/, js/), 跳过 data.js / data.json / assets/
        skip_items = {"data.js", "data.json", "assets"}
        for item in template_dir.iterdir():
            if item.name in skip_items:
                continue
            dst = self.output_dir / item.name
            try:
                if item.is_file():
                    shutil.copy2(item, dst)
                elif item.is_dir():
                    if dst.exists():
                        shutil.rmtree(dst)
                    shutil.copytree(item, dst)
                self.log(f"  复制: {item.name}")
            except Exception as e:
                self.log(f"  [warn] 复制失败 {item.name}: {e}")
        self.log("✓ 网站结构已生成")

    def _count_total_files(self) -> int:
        """快速统计主目录下符合条件的文件总数 (用于进度条)"""
        count = 0
        for dirpath, dirnames, filenames in os.walk(self.root):
            # 跳过忽略目录
            dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS and not d.startswith(".")]
            for name in filenames:
                ext = os.path.splitext(name)[1].lower().lstrip(".")
                if ext in SUPPORTED_EXTS and ext in self.ext_filter:
                    count += 1
        return count

    def _update_scan_progress(self):
        """根据已处理文件数更新扫描阶段进度 (3% ~ 70%)"""
        if self._total_files <= 0:
            return
        pct = 3 + int(67 * self._processed_files / self._total_files)
        if pct > 70:
            pct = 70
        self.set_progress(pct, f"扫描中... {self._processed_files}/{self._total_files}")

    def cancel(self):
        """请求停止扫描 (线程安全: 标志位写入)"""
        self._cancelled = True
        try:
            self.log("⚠ 收到停止请求, 正在停止扫描...")
        except Exception:
            pass

    def _is_cancelled(self) -> bool:
        """检查是否已请求停止"""
        return self._cancelled

    def log(self, msg: str):
        self.log_fn(msg)

    def set_progress(self, pct: int, text: str = ""):
        if self.progress_cb:
            self.progress_cb(min(100, max(0, int(pct))), text)

    def scan(self) -> dict:
        self.log(f"开始扫描: {self.root}")
        self.log(f"缩略图目录: {self.thumb_dir}")
        self.log(f"图片数量: {self.image_count}")
        self.log(f"文件类型: {', '.join(sorted(self.ext_filter))}")
        self.log(f"强制重扫: {'是' if self.force_rescan else '否 (增量)'}")

        start = time.time()
        self.files_flat = []
        self.current_paths = set()
        self.thumb_count = 0
        self.fail_count = 0
        self.skip_count = 0
        self.skipped_ext = 0
        self.new_count = 0
        self.cached_count = 0

        self.set_progress(1, "检测网站结构...")
        self._ensure_site_structure()

        self.set_progress(2, "统计文件总数...")
        self._total_files = self._count_total_files()
        self._processed_files = 0
        self.set_progress(3, f"开始扫描 (共 {self._total_files} 个文件)...")
        tree = self._scan_dir(self.root, depth=1, parent_id="")

        if self._is_cancelled():
            self.log("⏹ 扫描已停止, 不生成数据文件")
            self.set_progress(0, "已停止")
            return {"cancelled": True}

        self.set_progress(70, "清理已删除文件的缓存...")
        deleted_count = self.cache.sync_deleted(self.current_paths, self.thumb_dir, log=self.log)

        self.set_progress(80, "生成数据文件...")
        elapsed = time.time() - start

        # 统计后缀
        ext_counter = {}
        ext_category = {}
        for f in self.files_flat:
            ext = (f.get("ext") or "").lower()
            if not ext:
                continue
            ext_counter[ext] = ext_counter.get(ext, 0) + 1
            ext_category[ext] = f.get("category", "")
        exts = [
            {"ext": e, "count": c, "category": ext_category.get(e, "")}
            for e, c in sorted(ext_counter.items(), key=lambda kv: (-kv[1], kv[0]))
        ]

        data = {
            "version": "2.0",
            "root": self.root,
            "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
            "stats": {
                "folders": self._count_folders(tree),
                "files": len(self.files_flat),
                "thumbnails": self.thumb_count,
                "failed": self.fail_count,
                "skipped": self.skip_count,
                "deleted": deleted_count,
                "new": self.new_count,
                "cached": self.cached_count,
                "elapsed_sec": round(elapsed, 2),
            },
            "exts": exts,
            "tree": tree,
            "files": self.files_flat,
        }

        # 写出文件
        out_json = self.data_dir / "data.json"
        with open(out_json, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        try:
            shutil.copy2(out_json, self.root_data_json)
        except Exception:
            pass
        out_js = self.output_dir / "data.js"
        with open(out_js, "w", encoding="utf-8") as f:
            f.write("/* Auto-generated by scanner.py v3. Do not edit. */\n")
            f.write("window.__DOC_DATA = ")
            json.dump(data, f, ensure_ascii=False, indent=2)
            f.write(";\n")

        self.set_progress(95, "保存缓存...")
        self.cache.save()

        self.set_progress(100, "完成")
        self.log(f"\n========== 扫描完成 ==========")
        self.log(f"目录数:     {data['stats']['folders']}")
        self.log(f"文件数:     {data['stats']['files']}")
        self.log(f"新增:       {self.new_count}")
        self.log(f"命中缓存:   {self.cached_count}")
        self.log(f"删除清理:   {deleted_count}")
        self.log(f"缩略图数:   {data['stats']['thumbnails']}")
        self.log(f"失败数:     {data['stats']['failed']}")
        self.log(f"跳过:       {self.skip_count} (目录) + {self.skipped_ext} (扩展名)")
        self.log(f"耗时:       {elapsed:.2f}s")
        if exts:
            self.log(f"后缀种类:   {len(exts)} 种")
            for e in exts:
                self.log(f"  .{e['ext']:<8} {e['count']:>5} 个  ({e['category']})")
        self.log(f"数据文件:   {out_json}")
        return data

    def _count_folders(self, node_list):
        n = 0
        for node in node_list:
            if node.get("type") == "folder":
                n += 1
                n += self._count_folders(node.get("children", []))
        return n

    def _scan_dir(self, dir_path: str, depth: int, parent_id: str) -> list:
        try:
            entries = sorted(os.listdir(dir_path))
        except PermissionError:
            self.log(f"  [skip] 无访问权限: {dir_path}")
            return []
        except Exception as e:
            self.log(f"  [skip] 列目录失败 {dir_path}: {e}")
            return []

        nodes = []
        folders = []
        files = []
        for name in entries:
            full = os.path.join(dir_path, name)
            if os.path.isdir(full):
                if name in SKIP_DIRS or name.startswith("."):
                    self.skip_count += 1
                    continue
                folders.append(name)
            elif os.path.isfile(full):
                ext = os.path.splitext(name)[1].lower().lstrip(".")
                if ext in SUPPORTED_EXTS and ext in self.ext_filter:
                    files.append((name, ext))
                else:
                    self.skipped_ext += 1

        for name in folders:
            if self._is_cancelled():
                self.log(f"  [stop] 跳过目录: {name}")
                break
            full = os.path.join(dir_path, name)
            fid = safe_id(full)
            rel = os.path.relpath(full, self.root).replace("\\", "/")
            self.log(f"[D{depth}] 目录: {rel}")
            children = self._scan_dir(full, depth + 1, fid)
            nodes.append({
                "id": fid, "pid": parent_id, "type": "folder",
                "name": name, "path": full, "rel_path": rel,
                "depth": depth, "children": children,
            })

        for name, ext in files:
            if self._is_cancelled():
                self.log(f"  [stop] 跳过文件: {name}")
                break
            full = os.path.join(dir_path, name)
            fid = safe_id(full)
            rel = os.path.relpath(full, self.root).replace("\\", "/")
            self.log(f"[D{depth}] 文件: {rel}")

            # 更新进度
            self._processed_files += 1
            self._update_scan_progress()

            # 增量检查
            thumb_rel, is_new, all_thumbs = self._handle_file(full, ext, fid)
            size = os.path.getsize(full)
            file_node = {
                "id": fid, "pid": parent_id, "type": "file",
                "name": name, "ext": ext, "category": SUPPORTED_EXTS[ext],
                "path": full, "rel_path": rel, "depth": depth,
                "size": size, "size_text": human_size(size),
                "thumbnail": thumb_rel,
                "thumbnails": all_thumbs if all_thumbs else [thumb_rel] if thumb_rel else [],
            }
            nodes.append(file_node)
            self.files_flat.append(file_node)
            self.current_paths.add(full)
            if is_new:
                self.new_count += 1
            else:
                self.cached_count += 1

        return nodes

    def _handle_file(self, file_path: str, ext: str, fid: str) -> Tuple[str, bool, list]:
        """处理单个文件, 返回 (封面缩略图路径, 是否新增/变动, 全部缩略图路径列表)"""
        cat = SUPPORTED_EXTS.get(ext, "text")
        out_name = f"{fid}.png"
        out_path = self.thumb_dir / out_name
        thumb_rel = f"assets/thumbnails/{out_name}"

        def _collect_all_thumbs():
            """收集 fid 对应的所有缩略图 (封面 + _2, _3 ...)"""
            thumbs = [thumb_rel]
            for i in range(2, 11):
                multi_name = f"{fid}_{i}.png"
                if (self.thumb_dir / multi_name).exists():
                    thumbs.append(f"assets/thumbnails/{multi_name}")
                else:
                    break
            return thumbs

        # 检查缓存
        if not self.force_rescan and self.cache.is_processed(file_path):
            if out_path.exists():
                return thumb_rel, False, _collect_all_thumbs()
            # 缩略图被手动删除了, 重新生成
            self.log(f"  [regen] 缩略图缺失, 重新生成: {file_path}")

        # 生成缩略图
        ok = self._generate_thumbnail(file_path, ext, cat, out_path)
        if ok and out_path.exists():
            self.cache.put(file_path, out_name, ext, cat)
            self.thumb_count += 1
            return thumb_rel, True, _collect_all_thumbs()
        else:
            self.fail_count += 1
            try:
                gen_generic_thumbnail(file_path, ext, out_path)
            except Exception:
                pass
            thumb_rel2 = thumb_rel if out_path.exists() else ""
            if thumb_rel2:
                self.cache.put(file_path, out_name, ext, cat)
            return thumb_rel2, True, [thumb_rel2] if thumb_rel2 else []

    def _generate_thumbnail(self, file_path: str, ext: str, cat: str, out_path: Path) -> bool:
        try:
            if cat == "pdf":
                return gen_pdf_thumbnail(file_path, out_path, page_count=self.image_count)
            elif cat == "word":
                return gen_word_thumbnail(file_path, ext, out_path, page_count=self.image_count)
            elif cat == "ppt":
                return gen_ppt_thumbnail(file_path, ext, out_path, page_count=self.image_count)
            elif cat == "ebook":
                return gen_epub_thumbnail(file_path, out_path)
            elif cat == "text":
                return gen_text_thumbnail(file_path, out_path, ext=ext)
            else:
                return gen_generic_thumbnail(file_path, ext, out_path)
        except Exception as e:
            self.log(f"  [error] 缩略图异常 {file_path}: {e}")
            traceback.print_exc()
            return False


# ================================================================
# Tkinter GUI
# ================================================================
class ScannerGUI:
    TYPE_GROUPS = {
        "Documents":  {"label": "文档",     "exts": ["pdf", "doc", "docx", "rtf"]},
        "Slides":     {"label": "幻灯片",   "exts": ["ppt", "pptx"]},
        "Ebooks":     {"label": "电子书",   "exts": ["epub"]},
        "Text":       {"label": "文本/配置", "exts": ["txt", "md", "log", "csv", "json", "xml", "yaml", "yml", "ini", "cfg", "conf"]},
        "Code_Web":   {"label": "Web前端",   "exts": ["html", "htm", "css", "js", "jsx", "ts", "tsx", "vue", "svelte"]},
        "Code_Script":{"label": "脚本语言",  "exts": ["py", "rb", "pl", "php", "bat", "sh", "ps1", "r", "jl", "lua"]},
        "Code_Compile":{"label": "编译/系统", "exts": ["c", "h", "cpp", "cxx", "cc", "hpp", "cs", "java", "go", "rs", "kt", "swift", "sql"]},
        "Code_DotNet":{"label": ".NET 技术栈", "exts": ["cs", "vb", "aspx", "cshtml", "razor"]},
    }

    def __init__(self):
        import tkinter as tk
        from tkinter import ttk, filedialog, messagebox, scrolledtext

        self.tk = tk
        self.ttk = ttk
        self.filedialog = filedialog
        self.messagebox = messagebox
        self.scrolledtext = scrolledtext

        # 关键: 在创建 Tk 窗口前设置 AppUserModelID, 才能让任务栏/Alt-Tab
        #       显示本应用图标而非默认 Python 图标 (Windows 7+)
        set_app_user_model_id()

        self.root = tk.Tk()
        self.root.title("本地文档浏览器 - 扫描器")
        self.root.geometry("760x860")
        self.root.minsize(720, 780)
        # 设置窗口左上角图标 (多尺寸 .ico, 系统按需选择最合适尺寸)
        icon_file = get_icon_path()
        if icon_file:
            try:
                self.root.iconbitmap(default=icon_file)
                self.root.iconbitmap(icon_file)
            except Exception:
                pass

        self._build_ui()

    def _build_ui(self):
        tk = self.tk
        ttk = self.ttk
        frm = ttk.Frame(self.root, padding=10)
        frm.pack(fill=tk.BOTH, expand=True)

        # --- 目录选择 ---
        grp_dir = ttk.LabelFrame(frm, text="① 目录设置", padding=10)
        grp_dir.pack(fill=tk.X, pady=(0, 8))

        ttk.Label(grp_dir, text="扫描本地目录:").grid(row=0, column=0, sticky=tk.W, padx=(0, 6), pady=3)
        self.var_root = tk.StringVar()
        ttk.Entry(grp_dir, textvariable=self.var_root, width=55).grid(row=0, column=1, sticky=tk.EW, pady=3)
        ttk.Button(grp_dir, text="浏览…", command=self._choose_root).grid(row=0, column=2, padx=4, pady=3)

        ttk.Label(grp_dir, text="导出保存目录:").grid(row=1, column=0, sticky=tk.W, padx=(0, 6), pady=3)
        default_out = str(Path(__file__).resolve().parent.parent / "site")
        self.var_out = tk.StringVar(value=default_out)
        ttk.Entry(grp_dir, textvariable=self.var_out, width=55).grid(row=1, column=1, sticky=tk.EW, pady=3)
        ttk.Button(grp_dir, text="浏览…", command=self._choose_out).grid(row=1, column=2, padx=4, pady=3)
        grp_dir.columnconfigure(1, weight=1)

        # --- Web 服务器控制 ---
        grp_srv = ttk.LabelFrame(frm, text="  Web 服务", padding=10)
        grp_srv.pack(fill=tk.X, pady=(0, 8))

        self.var_port = tk.IntVar(value=8080)
        ttk.Label(grp_srv, text="端口:").grid(row=0, column=0, sticky=tk.W, padx=(0, 4))
        ttk.Spinbox(grp_srv, from_=1024, to=65535, textvariable=self.var_port, width=7).grid(row=0, column=1, sticky=tk.W, padx=(0, 12))

        self.lbl_srv_status = ttk.Label(grp_srv, text="● 未启动", foreground="#999")
        self.lbl_srv_status.grid(row=0, column=2, sticky=tk.W, padx=(0, 12))

        self.btn_srv_start = ttk.Button(grp_srv, text="启动服务", command=self._start_server, width=8)
        self.btn_srv_start.grid(row=0, column=3, padx=(0, 4))
        self.btn_srv_stop = ttk.Button(grp_srv, text="停止服务", command=self._stop_server, width=8, state="disabled")
        self.btn_srv_stop.grid(row=0, column=4, padx=(0, 4))
        self.btn_srv_open = ttk.Button(grp_srv, text="在浏览器中打开", command=self._open_browser, state="disabled")
        self.btn_srv_open.grid(row=0, column=5, padx=(0, 4))

        self.web_server = None  # WebServer 实例
        self._scanner = None    # 当前扫描器实例 (用于停止扫描)

        # --- 文件类型勾选 (按扩展名后缀显示) ---
        grp_type = ttk.LabelFrame(frm, text="② 扫描的文件类型 (按后缀勾选)", padding=10)
        grp_type.pack(fill=tk.X, pady=(0, 8))

        # 收集所有扩展名 (按优先级分组排序: 文档/电子书优先, 再代码, 最后其他)
        priority_order = [
            "pdf", "doc", "docx", "rtf",           # 文档
            "epub", "mobi", "azw3",                # 电子书
            "ppt", "pptx",                         # 幻灯片
            "txt", "md",                            # 常见文本
        ]
        all_exts = set()
        for info in self.TYPE_GROUPS.values():
            all_exts.update(info["exts"])
        # 优先扩展在前(保持指定顺序), 其余按字母序追加
        prioritized = [e for e in priority_order if e in all_exts]
        remaining = sorted(all_exts - set(prioritized), key=str.lower)
        self._all_exts = prioritized + remaining

        # 操作栏: 全选 / 全不选 / 反选
        bar = ttk.Frame(grp_type)
        bar.pack(fill=tk.X, pady=(0, 6))
        ttk.Button(bar, text="全选", width=6, command=lambda: self._toggle_all(True)).pack(side=tk.LEFT, padx=(0, 4))
        ttk.Button(bar, text="全不选", width=6, command=lambda: self._toggle_all(False)).pack(side=tk.LEFT, padx=(0, 4))
        ttk.Button(bar, text="反选", width=6, command=self._invert_selection).pack(side=tk.LEFT, padx=(0, 4))
        self.lbl_ext_count = ttk.Label(bar, text="", foreground="#555")
        self.lbl_ext_count.pack(side=tk.LEFT, padx=(12, 0))

        # 可滚动勾选区
        canvas_area = tk.Frame(grp_type, height=200, bd=1, relief=tk.SUNKEN)
        canvas_area.pack(fill=tk.X, pady=(0, 6))
        canvas_area.pack_propagate(False)

        self._canvas = tk.Canvas(canvas_area, highlightthickness=0)
        scrollbar = ttk.Scrollbar(canvas_area, orient=tk.VERTICAL, command=self._canvas.yview)
        self._canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self._canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        self._ext_frame = ttk.Frame(self._canvas)
        self._canvas_window = self._canvas.create_window((0, 0), window=self._ext_frame, anchor=tk.NW)

        # 绑定滚动 & 自适应宽度
        self._ext_frame.bind("<Configure>", self._on_ext_frame_configure)
        self._canvas.bind("<Configure>", self._on_canvas_configure)
        self._canvas.bind("<MouseWheel>", self._on_mousewheel)

        # 渲染勾选框
        self._ext_vars = {}  # ext -> BooleanVar
        self._render_ext_checkboxes()

        # 自定义添加后缀
        custom_frame = ttk.Frame(grp_type)
        custom_frame.pack(fill=tk.X, pady=(2, 0))
        ttk.Label(custom_frame, text="添加自定义后缀:").pack(side=tk.LEFT, padx=(0, 4))
        self.var_custom_ext = tk.StringVar()
        ttk.Entry(custom_frame, textvariable=self.var_custom_ext, width=12).pack(side=tk.LEFT, padx=(0, 4))
        ttk.Button(custom_frame, text="添加", command=self._add_custom_ext).pack(side=tk.LEFT, padx=(0, 6))
        ttk.Label(custom_frame, text="例: .pdf  .md  .csproj", foreground="#888").pack(side=tk.LEFT)

        self._update_ext_count()

        # --- 图片数量 ---
        grp_img = ttk.LabelFrame(frm, text="③ 缩略/预览图设置", padding=10)
        grp_img.pack(fill=tk.X, pady=(0, 8))

        ttk.Label(grp_img, text="提取页数 (1-10):").grid(row=0, column=0, sticky=tk.W, padx=(0, 10), pady=3)
        self.var_pages = tk.IntVar(value=1)
        spin = ttk.Spinbox(grp_img, from_=1, to=10, textvariable=self.var_pages, width=5)
        spin.grid(row=0, column=1, sticky=tk.W, pady=3)
        ttk.Label(grp_img, text="提示: 单页清晰, 多页占空间. 建议 ≤ 3 页，预估每张图片占用20~300KB",
                  foreground="#888").grid(row=0, column=2, sticky=tk.W, padx=10, pady=3)

        self.var_force = tk.BooleanVar(value=False)
        ttk.Checkbutton(grp_img, text="强制重新扫描 (忽略缓存)",
                        variable=self.var_force).grid(row=1, column=0, columnspan=3, sticky=tk.W, pady=(6, 0))

        # --- 操作按钮 ---
        grp_btn = ttk.Frame(frm)
        grp_btn.pack(fill=tk.X, pady=(0, 8))
        self.btn_scan = ttk.Button(grp_btn, text="开始扫描", command=self._start_scan)
        self.btn_scan.pack(side=tk.LEFT, padx=(0, 6))
        self.btn_stop = ttk.Button(grp_btn, text="停止扫描", command=self._stop_scan, state="disabled")
        self.btn_stop.pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(grp_btn, text="删除缓存", command=self._clear_cache).pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(grp_btn, text="打开导出目录", command=self._open_output).pack(side=tk.LEFT)

        # --- 进度条 ---
        self.progress = ttk.Progressbar(frm, mode="determinate", maximum=100)
        self.progress.pack(fill=tk.X, pady=(0, 6))
        self.lbl_progress = ttk.Label(frm, text="就绪")
        self.lbl_progress.pack(anchor=tk.W)

        # --- 日志区 ---
        grp_log = ttk.LabelFrame(frm, text="④ 日志", padding=4)
        grp_log.pack(fill=tk.BOTH, expand=True)
        self.txt_log = self.scrolledtext.ScrolledText(grp_log, height=14, wrap=tk.WORD,
                                                     font=("Consolas", 9))
        self.txt_log.pack(fill=tk.BOTH, expand=True)

        self._log("欢迎使用本地文档浏览器扫描器 v3")
        self._log("设置好后点击『开始扫描』即可")

    def _log(self, msg: str):
        ts = time.strftime("%H:%M:%S")
        self.txt_log.insert("end", f"[{ts}] {msg}\n")
        self.txt_log.see("end")
        self.root.update_idletasks()

    # ---- 扩展名勾选区: 渲染与交互 ----
    def _render_ext_checkboxes(self):
        """渲染所有扩展名勾选框 (显示为 .ext 格式)"""
        tk = self.tk
        ttk = self.ttk
        # 清空旧内容
        for child in self._ext_frame.winfo_children():
            child.destroy()
        self._ext_vars.clear()

        cols = 9  # 每行 9 个
        for idx, ext in enumerate(self._all_exts):
            row, col = divmod(idx, cols)
            var = tk.BooleanVar(value=True)
            self._ext_vars[ext] = var
            cb = ttk.Checkbutton(self._ext_frame, text=f".{ext}", variable=var,
                                 command=self._update_ext_count)
            cb.grid(row=row, column=col, sticky=tk.W, padx=(4, 10), pady=1)

    def _on_ext_frame_configure(self, event=None):
        """内容变化时更新滚动范围"""
        self._canvas.configure(scrollregion=self._canvas.bbox("all"))

    def _on_canvas_configure(self, event=None):
        """Canvas 大小变化时, 让内部 frame 宽度跟随"""
        self._canvas.itemconfig(self._canvas_window, width=event.width)

    def _on_mousewheel(self, event=None):
        """鼠标滚轮滚动"""
        delta = -1 if event.delta > 0 else 1
        self._canvas.yview_scroll(delta, "units")

    def _toggle_all(self, state: bool):
        """全选 / 全不选"""
        for var in self._ext_vars.values():
            var.set(state)
        self._update_ext_count()

    def _invert_selection(self):
        """反选"""
        for var in self._ext_vars.values():
            var.set(not var.get())
        self._update_ext_count()

    def _update_ext_count(self):
        """更新底部计数提示"""
        total = len(self._ext_vars)
        selected = sum(1 for v in self._ext_vars.values() if v.get())
        self.lbl_ext_count.config(text=f"已勾选 {selected} / {total} 个后缀")

    def _add_custom_ext(self):
        """添加自定义扩展名"""
        raw = self.var_custom_ext.get().strip().lstrip(".").lower()
        if not raw:
            self._log("⚠ 请输入扩展名, 例如 .pdf")
            return
        if not raw.replace(".", "").isalnum():
            self._log(f"⚠ 扩展名只能包含字母数字: {raw}")
            return
        if raw in self._ext_vars:
            # 已存在则勾选并提示
            if not self._ext_vars[raw].get():
                self._ext_vars[raw].set(True)
                self._update_ext_count()
                self._log(f"已勾选已有后缀: .{raw}")
            else:
                self._log(f"该后缀已存在且已勾选: .{raw}")
            return

        # 新增
        self._all_exts.append(raw)
        self._all_exts = sorted(self._all_exts, key=str.lower)
        self._render_ext_checkboxes()
        self._update_ext_count()
        self.var_custom_ext.set("")
        self._log(f"✓ 已添加自定义后缀: .{raw}")
        # 滚动到新项
        self.root.update_idletasks()
        self._canvas.yview_moveto(1.0)

    def _get_selected_exts(self) -> set:
        """获取已勾选的扩展名集合 (不含前导点)"""
        return {ext for ext, var in self._ext_vars.items() if var.get()}

    def _choose_root(self):
        d = self.filedialog.askdirectory(title="选择要浏览的文档主目录")
        if d:
            self.var_root.set(d)

    def _choose_out(self):
        d = self.filedialog.askdirectory(title="选择静态网站输出目录")
        if d:
            self.var_out.set(d)
            # 输出目录变化时, 如果服务器在运行则重启
            if self.web_server and self.web_server.running:
                self._log("输出目录已变更, 重启 Web 服务...")
                self._stop_server()
                self._start_server()

    # ---- Web 服务器控制 ----
    def _auto_start_server(self):
        """扫描完成后自动启动服务器"""
        if self.web_server and self.web_server.running:
            self._log("Web 服务已在运行")
            return
        out_dir = self.var_out.get().strip()
        if out_dir and os.path.exists(os.path.join(out_dir, "index.html")):
            self._start_server()
        else:
            self._log("输出目录无 index.html, 跳过自动启动")

    def _start_server(self):
        out_dir = self.var_out.get().strip()
        if not out_dir or not os.path.isdir(out_dir):
            self.messagebox.showwarning("提示", "输出目录不存在, 请先扫描生成数据")
            return
        # 检查 index.html
        if not os.path.exists(os.path.join(out_dir, "index.html")):
            self.messagebox.showwarning("提示", "输出目录中未找到 index.html, 请先扫描")
            return
        port = self.var_port.get()
        self.web_server = WebServer(out_dir, port=port, log_fn=self._log)
        ok, url = self.web_server.start()
        if ok:
            self.lbl_srv_status.config(text=f"● 运行中  {url}", foreground="#22c55e")
            self.btn_srv_start.config(state="disabled")
            self.btn_srv_stop.config(state="normal")
            self.btn_srv_open.config(state="normal")
            self._server_url = url
        else:
            self.messagebox.showerror("启动失败", url)

    def _stop_server(self):
        if self.web_server:
            self.web_server.stop()
        self.web_server = None
        self.lbl_srv_status.config(text="● 未启动", foreground="#999")
        self.btn_srv_start.config(state="normal")
        self.btn_srv_stop.config(state="disabled")
        self.btn_srv_open.config(state="disabled")

    def _open_browser(self):
        url = getattr(self, "_server_url", "") or f"http://127.0.0.1:{self.var_port.get()}"
        try:
            import webbrowser
            webbrowser.open(url)
            self._log(f"已打开浏览器: {url}")
        except Exception as e:
            self._log(f"打开浏览器失败: {e}")

    def _clear_cache(self):
        out = self.var_out.get().strip()
        if not out:
            return
        root = Path(out)
        # 需要删除的单独文件
        del_files = [
            root / "assets" / "data" / "scan_cache.json",
            root / "assets" / "data" / "data.json",
            root / "data.js"
        ]
        thumb_dir = root / "assets" / "thumbnails"

        msg = [
            "确定清空缓存？将删除：",
            "- assets/data/scan_cache.json",
            "- assets/data/data.json",
            "- data.js",
            "- assets/thumbnails/*.png"
        ]
        if not self.messagebox.askyesno("确认", "\n".join(msg)):
            return
        try:
            # 删除单个文件
            for f in del_files:
                if f.exists():
                    f.unlink()
                    self._log(f"已删除: {f}")
            # 删除缩略图目录下所有png
            if thumb_dir.is_dir():
                for png in thumb_dir.glob("*.png"):
                    png.unlink()
                    self._log(f"已删除缩略图: {png}")
        except Exception as e:
            self._log(f"删除缓存异常: {e}")


    def _open_output(self):
        out = self.var_out.get().strip()
        if out and os.path.isdir(out):
            try:
                if os.name == "nt":
                    os.startfile(out)
                else:
                    self._log(f"请手动打开: {out}")
            except Exception as e:
                self._log(f"打开失败: {e}")
        else:
            self.messagebox.showwarning("提示", "输出目录不存在")

    def _start_scan(self):
        root_dir = self.var_root.get().strip()
        out_dir = self.var_out.get().strip()
        if not root_dir or not os.path.isdir(root_dir):
            self.messagebox.showwarning("提示", "请选择有效的主目录")
            return
        if not out_dir:
            self.messagebox.showwarning("提示", "请选择输出目录")
            return
        exts = self._get_selected_exts()
        if not exts:
            self.messagebox.showwarning("提示", "请至少勾选一种文件类型")
            return

        self.btn_scan.config(state="disabled")
        self.btn_stop.config(state="normal")
        self.progress["value"] = 0
        self._scanner = None

        def run():
            try:
                scanner = Scanner(
                    root_dir, out_dir,
                    ext_filter=exts,
                    image_count=self.var_pages.get(),
                    force_rescan=self.var_force.get(),
                    log_fn=self._log,
                )
                scanner.progress_cb = self._set_progress
                self._scanner = scanner
                result = scanner.scan()
                if isinstance(result, dict) and result.get("cancelled"):
                    self._log("⏹ 扫描已停止")
                else:
                    self._log("✓ 扫描完成!")
                    # 自动启动 Web 服务
                    self.root.after(100, self._auto_start_server)
                    self.messagebox.showinfo("完成", f"扫描完成!\nWeb 服务已自动启动\n点击「在浏览器中打开」即可使用")
            except Exception as e:
                self._log(f"[error] 扫描异常: {e}")
                traceback.print_exc()
                self.messagebox.showerror("错误", f"扫描异常:\n{e}")
            finally:
                self._scanner = None
                self.btn_scan.config(state="normal")
                self.btn_stop.config(state="disabled")

        threading.Thread(target=run, daemon=True).start()

    def _stop_scan(self):
        if self._scanner:
            self._scanner.cancel()
            self._log("正在停止扫描...")
            self.btn_stop.config(state="disabled")
        else:
            self.messagebox.showinfo("提示", "当前没有正在进行的扫描")

    def _set_progress(self, pct: int, text: str):
        self.progress["value"] = pct
        self.lbl_progress.config(text=text)
        self.root.update_idletasks()

    def run(self):
        self.root.mainloop()


# ================================================================
# 入口 (支持命令行与 GUI)
# ================================================================
def main():
    args = sys.argv[1:]

    # 如果带参数, 走命令行模式
    if args:
        run_cli(args)
        return

    # 否则启动 GUI
    try:
        gui = ScannerGUI()
        gui.run()
    except Exception as e:
        print(f"GUI 启动失败: {e}")
        print("可使用命令行模式: python scanner.py --root <path> --out <path> --ext pdf,txt --pages 2")
        input("按回车退出...")


def run_cli(args):
    root_dir = ""
    output_dir = ""
    pages = 1
    exts_filter = None
    force = False

    i = 0
    while i < len(args):
        a = args[i]
        if a == "--root" and i + 1 < len(args):
            root_dir = args[i + 1]; i += 2
        elif a == "--out" and i + 1 < len(args):
            output_dir = args[i + 1]; i += 2
        elif a == "--pages" and i + 1 < len(args):
            try:
                pages = int(args[i + 1])
            except Exception:
                pass
            i += 2
        elif a == "--ext" and i + 1 < len(args):
            exts_filter = set(e.strip().lower().lstrip(".") for e in args[i + 1].split(","))
            i += 2
        elif a == "--force":
            force = True; i += 1
        else:
            i += 1

    if not root_dir:
        print("缺少 --root 参数")
        sys.exit(1)
    if not os.path.isdir(root_dir):
        print(f"目录不存在: {root_dir}")
        sys.exit(1)
    if not output_dir:
        output_dir = str(Path(__file__).resolve().parent.parent / "site")

    scanner = Scanner(root_dir, output_dir,
                      ext_filter=exts_filter,
                      image_count=pages,
                      force_rescan=force)
    scanner.scan()


if __name__ == "__main__":
    main()
